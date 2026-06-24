"""
ADAM slidedeck skill.

Creates local PowerPoint (.pptx) artifacts from structured slide plans or
markdown-style outline content. This skill does not verify claims, send files,
upload files, or access the network.
"""

from __future__ import annotations

import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

SKILL_NAME = "slidedeck"
SUPPORTED_ACTIONS = {"create"}
SUPPORTED_THEMES = {"governance", "light", "dark", "simple"}
MAX_SLIDES_HARD_CAP = 80
DEFAULT_MAX_SLIDES = 40
MAX_TITLE_CHARS = 180
MAX_BULLET_CHARS = 260
MAX_NOTES_CHARS = 3000


def _fail(action: str, error_class: str, error_message: str) -> Dict[str, Any]:
    return {
        "ok": False,
        "status": "failed",
        "skill": SKILL_NAME,
        "action": action,
        "error_class": error_class,
        "error_message": error_message,
    }


def _clean_text(value: Any, max_len: int = 1000) -> str:
    if value is None:
        return ""
    if not isinstance(value, str):
        value = str(value)
    value = value.replace("\x00", "").strip()
    value = re.sub(r"[ \t]+", " ", value)
    return value[:max_len].strip()


def _coerce_bool(value: Any, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        v = value.strip().lower()
        if v in {"true", "1", "yes", "y", "on"}:
            return True
        if v in {"false", "0", "no", "n", "off"}:
            return False
    return default


def _coerce_int(value: Any, default: int, low: int, high: int) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        parsed = default
    return max(low, min(high, parsed))


def _safe_filename(name: str, default: str = "adam_slidedeck") -> str:
    name = _clean_text(name, 140).lower()
    name = re.sub(r"[^a-z0-9._ -]", "", name)
    name = re.sub(r"[\s-]+", "_", name).strip("._")
    if not name:
        name = default
    if not name.endswith(".pptx"):
        name += ".pptx"
    return name


def _artifact_dir(context: Dict[str, Any]) -> Path:
    """
    Resolve the artifact directory for this skill invocation.

    The ADAM skill runtime passes the session's artifacts directory
    via context["artifacts_root"] (see adam/skills_runtime/runtime.py).
    That is the canonical key; every skill that writes a real artifact
    should resolve to it so that audit, GUI, and continuity all see
    the file in the expected place.

    Fallback order:
      1. context["artifacts_root"]  -- the runtime-provided canonical
         directory. Always present in production calls.
      2. context["artifact_dir"], context["artifacts_dir"],
         context["output_dir"] -- legacy/test keys, kept for
         compatibility with debug harnesses that may use them.
      3. context["log_dir"]/artifacts -- another legacy form.
      4. ./artifacts/slidedeck -- last resort for ad-hoc unit tests
         only; should never fire in a real session because (1) is
         always set by the runtime.
    """
    candidates = [
        context.get("artifacts_root"),
        context.get("artifact_dir"),
        context.get("artifacts_dir"),
        context.get("output_dir"),
        context.get("log_dir") and Path(str(context.get("log_dir"))) / "artifacts",
    ]
    for c in candidates:
        if not c:
            continue
        try:
            p = Path(str(c)).expanduser().resolve()
            p.mkdir(parents=True, exist_ok=True)
            return p
        except Exception:
            continue
    p = Path("artifacts") / "slidedeck"
    p.mkdir(parents=True, exist_ok=True)
    return p.resolve()


def _sanitize_output_path(output_dir: Path, filename: str) -> Path:
    safe = _safe_filename(filename)
    path = (output_dir / safe).resolve()
    if output_dir.resolve() not in path.parents and path != output_dir.resolve():
        raise ValueError("resolved output path escapes artifact directory")
    return path


def _listify(value: Any, max_items: int = 12, max_len: int = MAX_BULLET_CHARS) -> List[str]:
    if value is None:
        return []
    if isinstance(value, list):
        raw = value
    elif isinstance(value, tuple):
        raw = list(value)
    elif isinstance(value, str):
        raw = [line.strip(" -•\t") for line in value.splitlines() if line.strip()]
    else:
        raw = [str(value)]
    out: List[str] = []
    for item in raw:
        text = _clean_text(item, max_len)
        if text:
            out.append(text)
        if len(out) >= max_items:
            break
    return out


def _parse_markdown_content(content: str, max_slides: int) -> List[Dict[str, Any]]:
    """Parse a simple markdown outline into slides."""
    content = _clean_text(content, 100000)
    if not content:
        return []

    slides: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None

    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("#"):
            heading = line.lstrip("#").strip()
            if current:
                slides.append(current)
                if len(slides) >= max_slides:
                    return slides
            current = {"title": heading, "layout": "bullets", "bullets": []}
        elif line.startswith(("- ", "* ", "• ")):
            if current is None:
                current = {"title": "Key Points", "layout": "bullets", "bullets": []}
            current.setdefault("bullets", []).append(line[2:].strip()[:MAX_BULLET_CHARS])
        else:
            if current is None:
                current = {"title": line[:MAX_TITLE_CHARS], "layout": "bullets", "bullets": []}
            else:
                current.setdefault("bullets", []).append(line[:MAX_BULLET_CHARS])

    if current and len(slides) < max_slides:
        slides.append(current)
    return slides[:max_slides]


def _normalize_slide(raw: Any) -> Dict[str, Any]:
    if not isinstance(raw, dict):
        return {"title": _clean_text(raw, MAX_TITLE_CHARS) or "Slide", "layout": "bullets", "bullets": []}

    layout = _clean_text(raw.get("layout"), 32).lower() or "bullets"
    if layout not in {"title", "bullets", "two_column", "quote", "section", "closing"}:
        layout = "bullets"

    slide = {
        "title": _clean_text(raw.get("title"), MAX_TITLE_CHARS) or "Slide",
        "subtitle": _clean_text(raw.get("subtitle"), 300),
        "layout": layout,
        "bullets": _listify(raw.get("bullets"), max_items=10),
        "left_title": _clean_text(raw.get("left_title"), 120),
        "left_bullets": _listify(raw.get("left_bullets"), max_items=8),
        "right_title": _clean_text(raw.get("right_title"), 120),
        "right_bullets": _listify(raw.get("right_bullets"), max_items=8),
        "quote": _clean_text(raw.get("quote"), 600),
        "speaker_notes": _clean_text(raw.get("speaker_notes"), MAX_NOTES_CHARS),
    }
    return slide


def _normalize_slides(args: Dict[str, Any], max_slides: int) -> List[Dict[str, Any]]:
    raw_slides = args.get("slides")
    slides: List[Dict[str, Any]] = []
    if isinstance(raw_slides, list) and raw_slides:
        slides = [_normalize_slide(s) for s in raw_slides[:max_slides]]
    else:
        slides = [_normalize_slide(s) for s in _parse_markdown_content(args.get("content", ""), max_slides)]
    return slides[:max_slides]


def _theme(theme_name: str) -> Dict[str, str]:
    themes = {
        "governance": {
            "bg": "071014", "panel": "0F2027", "text": "E8F7F5", "muted": "8AA6AA",
            "accent": "6CF5E6", "accent2": "83F28F", "warn": "F4C76B",
        },
        "dark": {
            "bg": "111827", "panel": "1F2937", "text": "F9FAFB", "muted": "CBD5E1",
            "accent": "60A5FA", "accent2": "A78BFA", "warn": "FBBF24",
        },
        "light": {
            "bg": "F8FAFC", "panel": "FFFFFF", "text": "0F172A", "muted": "475569",
            "accent": "0EA5E9", "accent2": "16A34A", "warn": "B45309",
        },
        "simple": {
            "bg": "FFFFFF", "panel": "FFFFFF", "text": "111827", "muted": "4B5563",
            "accent": "2563EB", "accent2": "059669", "warn": "92400E",
        },
    }
    return themes.get(theme_name, themes["governance"])


def _add_textbox(slide: Any, x: float, y: float, w: float, h: float, text: str, size: int, color: str,
                 bold: bool = False, font: str = "Aptos", valign: Optional[str] = None) -> Any:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    if valign:
        try:
            from pptx.enum.text import MSO_VERTICAL_ANCHOR
            frame.vertical_anchor = getattr(MSO_VERTICAL_ANCHOR, valign)
        except Exception:
            pass
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_bullets(slide: Any, x: float, y: float, w: float, h: float, bullets: List[str], color: str,
                 size: int = 24, font: str = "Aptos") -> Any:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, b in enumerate(bullets or [" "]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = font
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor.from_string(color)
    return box


def _add_notes(slide: Any, notes: str) -> None:
    if not notes:
        return
    try:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes
    except Exception:
        # Speaker notes are helpful, but not load-bearing.
        pass


def _add_footer(slide: Any, footer: str, slide_no: int, total: int, theme: Dict[str, str], wide: bool) -> None:
    if not footer:
        footer = "ADAM"
    width = 13.333 if wide else 10
    height = 7.5 if wide else 7.5
    _add_textbox(slide, 0.55, height - 0.42, width - 1.1, 0.22, footer, 8, theme["muted"])
    _add_textbox(slide, width - 1.15, height - 0.42, 0.7, 0.22, f"{slide_no}/{total}", 8, theme["muted"])


def _set_background(slide: Any, color: str) -> None:
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def _add_accent_bar(slide: Any, theme: Dict[str, str], wide: bool) -> None:
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    width = 13.333 if wide else 10
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(width), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(theme["accent"])
    shape.line.fill.background()


def _create_presentation(args: Dict[str, Any], output_path: Path) -> Tuple[int, List[str]]:
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
    except Exception as e:
        raise RuntimeError(
            "python-pptx is required for the slidedeck skill. Install with: pip install python-pptx"
        ) from e

    title = _clean_text(args.get("title"), MAX_TITLE_CHARS)
    subtitle = _clean_text(args.get("subtitle"), 400)
    footer = _clean_text(args.get("footer"), 160)
    author = _clean_text(args.get("author"), 120)
    theme_name = _clean_text(args.get("theme"), 32).lower() or "governance"
    if theme_name not in SUPPORTED_THEMES:
        theme_name = "governance"
    aspect_ratio = _clean_text(args.get("aspect_ratio"), 32).lower() or "wide"
    wide = aspect_ratio != "standard"
    max_slides = _coerce_int(args.get("max_slides", DEFAULT_MAX_SLIDES), DEFAULT_MAX_SLIDES, 1, MAX_SLIDES_HARD_CAP)
    include_agenda = _coerce_bool(args.get("include_agenda"), False)
    include_section_dividers = _coerce_bool(args.get("include_section_dividers"), True)

    slide_defs = _normalize_slides(args, max_slides)
    if not slide_defs:
        slide_defs = [{"title": "Overview", "layout": "bullets", "bullets": ["No slide details were provided."]}]

    prs = Presentation()
    if wide:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

    try:
        prs.core_properties.title = title
        prs.core_properties.author = author or "ADAM"
        prs.core_properties.subject = "ADAM-generated governed slide deck"
        prs.core_properties.comments = "Created by ADAM slidedeck skill. Claims should be verified by Truthseeker when factual support is required."
    except Exception:
        pass

    colors = _theme(theme_name)
    blank_layout = prs.slide_layouts[6]
    generated_titles: List[str] = []

    def new_slide() -> Any:
        s = prs.slides.add_slide(blank_layout)
        _set_background(s, colors["bg"])
        _add_accent_bar(s, colors, wide)
        return s

    # Title slide
    slide = new_slide()
    _add_textbox(slide, 0.8, 1.65, 11.8 if wide else 8.4, 1.1, title, 42, colors["text"], bold=True)
    if subtitle:
        _add_textbox(slide, 0.85, 2.85, 10.9 if wide else 8.1, 0.75, subtitle, 21, colors["muted"])
    _add_textbox(slide, 0.85, 6.75, 8.0, 0.25, "Generated by ADAM Governance Core", 10, colors["accent"])
    generated_titles.append(title)

    visible_slide_defs = [s for s in slide_defs if include_section_dividers or s.get("layout") != "section"]
    if include_agenda and visible_slide_defs:
        slide = new_slide()
        _add_textbox(slide, 0.7, 0.55, 11.8 if wide else 8.5, 0.55, "Agenda", 30, colors["text"], bold=True)
        agenda_items = [s.get("title", "Slide") for s in visible_slide_defs[:9]]
        _add_bullets(slide, 1.0, 1.45, 10.8 if wide else 8.0, 4.8, agenda_items, colors["text"], size=23)
        generated_titles.append("Agenda")

    for sdef in visible_slide_defs:
        layout = sdef.get("layout", "bullets")
        slide = new_slide()
        title_text = sdef.get("title", "Slide")
        generated_titles.append(title_text)

        if layout == "section":
            _add_textbox(slide, 0.8, 2.45, 11.8 if wide else 8.4, 0.9, title_text, 40, colors["accent"], bold=True)
            if sdef.get("subtitle"):
                _add_textbox(slide, 0.85, 3.45, 10.8 if wide else 8.0, 0.7, sdef["subtitle"], 20, colors["muted"])
        elif layout == "quote":
            _add_textbox(slide, 0.75, 0.6, 11.8 if wide else 8.4, 0.5, title_text, 27, colors["text"], bold=True)
            quote = sdef.get("quote") or " ".join(sdef.get("bullets", [])) or sdef.get("subtitle", "")
            _add_textbox(slide, 1.25, 2.1, 10.7 if wide else 7.7, 1.8, f"“{quote}”", 30, colors["accent"], bold=False, valign="MIDDLE")
        elif layout == "two_column":
            _add_textbox(slide, 0.65, 0.5, 12.0 if wide else 8.7, 0.55, title_text, 28, colors["text"], bold=True)
            # Column panels
            for x, w in ((0.75, 5.75 if wide else 4.15), (6.85 if wide else 5.1, 5.75 if wide else 4.15)):
                panel = slide.shapes.add_shape(1, Inches(x), Inches(1.35), Inches(w), Inches(5.3))
                panel.fill.solid()
                panel.fill.fore_color.rgb = RGBColor.from_string(colors["panel"])
                panel.line.color.rgb = RGBColor.from_string(colors["accent"])
            _add_textbox(slide, 1.05, 1.6, 5.1 if wide else 3.6, 0.4, sdef.get("left_title") or "First", 20, colors["accent"], bold=True)
            _add_bullets(slide, 1.05, 2.25, 5.0 if wide else 3.6, 3.9, sdef.get("left_bullets", []), colors["text"], size=18)
            _add_textbox(slide, 7.15 if wide else 5.4, 1.6, 5.1 if wide else 3.6, 0.4, sdef.get("right_title") or "Second", 20, colors["accent2"], bold=True)
            _add_bullets(slide, 7.15 if wide else 5.4, 2.25, 5.0 if wide else 3.6, 3.9, sdef.get("right_bullets", []), colors["text"], size=18)
        elif layout == "closing":
            _add_textbox(slide, 0.8, 1.75, 11.8 if wide else 8.4, 0.8, title_text, 38, colors["text"], bold=True)
            if sdef.get("subtitle"):
                _add_textbox(slide, 0.85, 2.8, 10.8 if wide else 8.0, 0.6, sdef["subtitle"], 22, colors["muted"])
            if sdef.get("bullets"):
                _add_bullets(slide, 1.0, 3.65, 10.5 if wide else 7.8, 2.1, sdef.get("bullets", []), colors["text"], size=21)
        else:
            _add_textbox(slide, 0.65, 0.5, 12.0 if wide else 8.7, 0.55, title_text, 28, colors["text"], bold=True)
            if sdef.get("subtitle"):
                _add_textbox(slide, 0.7, 1.05, 11.2 if wide else 8.3, 0.4, sdef["subtitle"], 15, colors["muted"])
            _add_bullets(slide, 0.95, 1.65, 11.0 if wide else 8.0, 4.8, sdef.get("bullets", []), colors["text"], size=22)

        _add_notes(slide, sdef.get("speaker_notes", ""))

    total = len(prs.slides)
    for idx, sl in enumerate(prs.slides, start=1):
        _add_footer(sl, footer, idx, total, colors, wide)

    prs.save(output_path)
    return total, generated_titles


def handle(action: str, args: Dict[str, Any], context: Dict[str, Any]) -> Dict[str, Any]:
    """
    ADAM skill handler entry point.

    Expected signature:
        handle(action: str, args: dict, context: dict) -> dict
    """
    if action not in SUPPORTED_ACTIONS:
        return _fail(
            action=action,
            error_class="disallowed_action",
            error_message=f"Action '{action}' is unrecognized. Supported actions: ['create'].",
        )

    if not isinstance(args, dict):
        return _fail(action, "invalid_args", "args must be an object/dictionary.")
    if not isinstance(context, dict):
        context = {}

    title = _clean_text(args.get("title"), MAX_TITLE_CHARS)
    if not title:
        return _fail(
            action=action,
            error_class="missing_required_args",
            error_message="The 'title' argument is required and must be a non-empty string.",
        )

    if not args.get("slides") and not _clean_text(args.get("content"), 1000):
        return _fail(
            action=action,
            error_class="missing_required_content",
            error_message="Provide either 'slides' as an array or 'content' as a markdown-style outline.",
        )

    try:
        out_dir = _artifact_dir(context)
        requested_name = args.get("output_filename") or title
        output_path = _sanitize_output_path(out_dir, str(requested_name))
        if output_path.exists():
            stem = output_path.stem
            suffix = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = output_path.with_name(f"{stem}_{suffix}.pptx")

        slide_count, slide_titles = _create_presentation(args, output_path)
    except Exception as e:
        return _fail(
            action=action,
            error_class="slidedeck_creation_failed",
            error_message=f"Slide deck creation failed: {type(e).__name__}: {e}",
        )

    artifact_id = f"SLIDEDECK-{datetime.now().strftime('%Y%m%d-%H%M%S')}"
    return {
        "ok": True,
        "status": "success",
        "skill": SKILL_NAME,
        "action": action,
        "artifact_id": artifact_id,
        "path": str(output_path),
        "filename": output_path.name,
        "format": "pptx",
        "slide_count": slide_count,
        "slide_titles": slide_titles,
        "audit_meta": {
            "io_operation": "local_file_write",
            "write_access_asserted": True,
            "external_network_access_asserted": False,
            "human_approval_required": False,
            "truthseeker_followup_required": True,
        },
        "note": "Slide deck created as a local artifact. This skill does not verify claims or distribute the file.",
    }
